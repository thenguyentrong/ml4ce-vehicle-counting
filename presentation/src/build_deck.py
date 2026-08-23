"""Build the ML4CE Topic 2 presentation in the CTRL design language.

Same grid, type and components as the CTRL Summer School decks
(DesignSystem/src/build_gh_pptx.py): white paper, ink type, Montserrat + JetBrains Mono,
mono eyebrows with wide tracking, a hairline footer and a diamond page number.
No CTRL or partner logos - this is coursework, not a CTRL event.

20 slides. Speaker notes carry the timing.
"""
import json
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x11, 0x11, 0x11)
MUT = RGBColor(0x66, 0x66, 0x66)
SLATE = RGBColor(0x4A, 0x55, 0x68)
LINE = RGBColor(0xE4, 0xE4, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

F_BLACK = "Montserrat Black"
F_LIGHT = "Montserrat Light"
F_REG = "Montserrat"
F_MONO = "JetBrains Mono"

FIGS = os.environ["DECK_FIGS"]
OUT = os.environ["DECK_OUT"]

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]
ML = Inches(0.667)
MR = Inches(13.333 - 0.667)
CW = Inches(13.333 - 2 * 0.667)
FOOT = "ML4CE · SEMESTER PROJECT · TOPIC 2"

_PAGE = [1]


def autopage():
    _PAGE[0] += 1
    return _PAGE[0]


def track(run, hundredths):
    run.font._rPr.set("spc", str(hundredths))


def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def runs(para, parts, line_spacing=None, space_after=None, align=None):
    if line_spacing:
        para.line_spacing = line_spacing
    if space_after is not None:
        para.space_after = space_after
    if align:
        para.alignment = align
    for text, font, size, color, bold, spc in parts:
        r = para.add_run()
        r.text = text
        r.font.name = font
        r.font.size = size
        r.font.color.rgb = color
        r.font.bold = bold
        if spc:
            track(r, spc)
    return para


def rect(slide, x, y, w, h, color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def hairline(slide, y, x=ML, w=CW, color=LINE, weight=9525):
    return rect(slide, x, y, w, Emu(weight), color)


def eyebrow(slide, text, y=Inches(0.52), color=SLATE, x=ML, w=CW):
    _, tf = tb(slide, x, y, w, Inches(0.3))
    runs(tf.paragraphs[0], [(text.upper(), F_MONO, Pt(11), color, False, 150)])


def diamond(slide, num, cx=Inches(12.50), cy=Inches(6.90), size=Inches(0.30)):
    sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, cx, cy, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = INK
    sh.line.width = Pt(1.1)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    runs(tf.paragraphs[0], [(str(num), F_MONO, Pt(8.5), INK, False, 0)], align=PP_ALIGN.CENTER)


def footer(slide, left_extra=""):
    num = autopage()
    hairline(slide, Inches(6.72))
    _, tf = tb(slide, ML, Inches(6.84), Inches(9.5), Inches(0.3))
    parts = [("VEHICLE DETECTION & COUNTING", F_MONO, Pt(8), INK, False, 110),
             ("   ·   " + FOOT, F_MONO, Pt(8), MUT, False, 110)]
    if left_extra:
        parts.append(("   ·   " + left_extra.upper(), F_MONO, Pt(8), MUT, False, 110))
    runs(tf.paragraphs[0], parts)
    diamond(slide, num)


def content_header(slide, kicker, title, sub):
    eyebrow(slide, kicker)
    _, tf = tb(slide, ML, Inches(0.88), CW, Inches(0.55))
    runs(tf.paragraphs[0], [(title, F_LIGHT, Pt(26), INK, False, 0)])
    _, tf = tb(slide, ML, Inches(1.42), CW, Inches(0.4))
    runs(tf.paragraphs[0], [(sub, F_REG, Pt(15), INK, True, 0)])
    rect(slide, ML, Inches(1.88), CW, Emu(12700), INK)


def statement(slide, kicker, parts, sub, extra_mono=None):
    eyebrow(slide, kicker, y=Inches(1.30))
    _, tf = tb(slide, ML, Inches(2.05), Inches(10.9), Inches(3.4))
    runs(tf.paragraphs[0], parts, line_spacing=1.45)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(16)
    runs(p2, [(sub, F_LIGHT, Pt(16), MUT, False, 0)], line_spacing=1.4)
    if extra_mono:
        _, tf2 = tb(slide, ML, Inches(5.55), Inches(11), Inches(0.4))
        runs(tf2.paragraphs[0], [(extra_mono.upper(), F_MONO, Pt(11), SLATE, False, 130)])
    footer(slide)


def caption(slide, text, x, y, w=Inches(6.5)):
    _, tf = tb(slide, x, y, w, Inches(0.3))
    runs(tf.paragraphs[0], [(text, F_LIGHT, Pt(9.5), MUT, False, 0)])
    for run in tf.paragraphs[0].runs:
        run.font.italic = True


def col_label(slide, x, y, text, w=Inches(5)):
    _, tf = tb(slide, x, y, w, Inches(0.28))
    runs(tf.paragraphs[0], [(text.upper(), F_MONO, Pt(10), SLATE, False, 150)])


def chip_flow(slide, y, items, x=ML, w=CW, chip_h=Inches(0.95)):
    n = len(items)
    arrow_w = Inches(0.42)
    chip_w = Emu(int((w - arrow_w * (n - 1)) / n))
    xx = x
    for i, (label, sub) in enumerate(items):
        rect(slide, Emu(int(xx)), y, chip_w, chip_h, WHITE, line_color=INK)
        _, tf = tb(slide, Emu(int(xx + Inches(0.12))), Emu(int(y + Inches(0.15))),
                   Emu(int(chip_w - Inches(0.24))), Emu(int(chip_h - Inches(0.2))))
        runs(tf.paragraphs[0], [(label, F_MONO, Pt(10.5), INK, True, 60)], space_after=Pt(2))
        if sub:
            p = tf.add_paragraph()
            runs(p, [(sub, F_LIGHT, Pt(9.5), MUT, False, 0)], line_spacing=1.1)
        if i < n - 1:
            _, tf = tb(slide, Emu(int(xx + chip_w)), Emu(int(y + chip_h / 2 - Inches(0.15))),
                       arrow_w, Inches(0.3))
            runs(tf.paragraphs[0], [("→", F_MONO, Pt(13), MUT, False, 0)], align=PP_ALIGN.CENTER)
        xx += chip_w + arrow_w


def pic_fit(slide, name, x, y, w, h, align="center"):
    """Place an image inside the box without cropping it."""
    src = os.path.join(FIGS, name)
    img = Image.open(src)
    ar = img.width / img.height
    if w / h > ar:
        hh, ww = h, Emu(int(h * ar))
    else:
        ww, hh = w, Emu(int(w / ar))
    ox = x if align == "left" else Emu(int(x + (w - ww) / 2))
    oy = Emu(int(y + (h - hh) / 2))
    return slide.shapes.add_picture(src, ox, oy, width=ww, height=hh)


def rows(slide, x, y, items, w=Inches(11.6), gap=0.84, num_font=F_BLACK):
    """Numbered rows with a hairline under each - the deck's list component."""
    yy = y
    for num, head, body in items:
        _, tf = tb(slide, x, Inches(yy), w, Inches(0.8))
        runs(tf.paragraphs[0], [(num + "  ", num_font, Pt(14), INK, False, 0),
                                (head, F_REG, Pt(13.5), INK, True, 0)], space_after=Pt(2))
        p = tf.add_paragraph()
        runs(p, [(body, F_LIGHT, Pt(11.5), MUT, False, 0)], line_spacing=1.2)
        hairline(slide, Inches(yy + gap - 0.10), x=x, w=w)
        yy += gap
    return yy


def table(slide, x, y, cols, data, row_h=0.42, head=True, bold_rows=(), w=None):
    """Rule-based table: mono header, hairline per row, no fills and no borders."""
    widths = [c[1] for c in cols]
    total = sum(widths)
    w = w or CW
    xs, acc = [], 0.0
    for cw in widths:
        xs.append(x + Emu(int(w * acc / total)))
        acc += cw

    yy = y
    if head:
        for xi, (label, _) in zip(xs, cols):
            _, tf = tb(slide, xi, Inches(yy), Inches(3.4), Inches(0.3))
            runs(tf.paragraphs[0], [(label.upper(), F_MONO, Pt(9), SLATE, False, 140)])
        yy += 0.34
        rect(slide, x, Inches(yy), Emu(int(w)), Emu(12700), INK)
        yy += 0.12

    for i, row in enumerate(data):
        bold = i in bold_rows
        for j, (xi, cell) in enumerate(zip(xs, row)):
            _, tf = tb(slide, xi, Inches(yy), Emu(int(w * widths[j] / total - Inches(0.12))),
                       Inches(row_h))
            font = F_MONO if j and _is_num(cell) else F_REG
            runs(tf.paragraphs[0], [(cell, font, Pt(11.5) if bold else Pt(11),
                                     INK if bold else (INK if j == 0 else MUT), bold, 0)],
                 line_spacing=1.1)
        yy += row_h
        hairline(slide, Inches(yy - 0.08), x=x, w=Emu(int(w)))
    return yy


def _is_num(text):
    return any(ch.isdigit() for ch in text) and len(text) < 24


def notes(slide, text):
    """Write the speaker note into the slide's notes pane (View > Notes in PowerPoint)."""
    slide.notes_slide.notes_text_frame.text = text


def new():
    return prs.slides.add_slide(BLANK)


# ============================================================ 01 · Title
s = new()
eyebrow(s, "ML4CE · Semester project · Topic 2", y=Inches(0.52))
_, tf = tb(s, Inches(7.3), Inches(0.50), Inches(5.35), Inches(0.3))
runs(tf.paragraphs[0], [("RWTH AACHEN UNIVERSITY", F_MONO, Pt(11), SLATE, False, 160)],
     align=PP_ALIGN.RIGHT)
_, tf = tb(s, ML, Inches(2.05), Inches(12.2), Inches(2.4))
runs(tf.paragraphs[0], [("VEHICLE", F_BLACK, Pt(88), INK, False, 0)], line_spacing=0.95)
p = tf.add_paragraph()
runs(p, [("COUNTING", F_BLACK, Pt(88), INK, False, 0)], line_spacing=0.95)
_, tf = tb(s, ML, Inches(4.60), Inches(10.2), Inches(1.0))
runs(tf.paragraphs[0], [
    ("A detector built from scratch", F_REG, Pt(17), INK, True, 0),
    (", and a counting pipeline for traffic video.", F_LIGHT, Pt(17), INK, False, 0),
], line_spacing=1.45)
_, tf = tb(s, ML, Inches(5.95), Inches(12), Inches(0.35))
runs(tf.paragraphs[0], [("THE VINH NGUYEN TRONG   ·   ALI AWADA   ·   REXHEP AZEMI",
                         F_MONO, Pt(11.5), INK, False, 120)])
hairline(s, Inches(6.72))
_, tf = tb(s, ML, Inches(6.86), Inches(8), Inches(0.3))
runs(tf.paragraphs[0], [("PART 1 ", F_MONO, Pt(8.5), INK, False, 120),
                        ("DETECTOR FROM SCRATCH   ·   ", F_MONO, Pt(8.5), MUT, False, 120),
                        ("PART 2 ", F_MONO, Pt(8.5), INK, False, 120),
                        ("COUNTING IN VIDEO", F_MONO, Pt(8.5), MUT, False, 120)])
notes(s, "~30 SEC. - Question first, then the names. Short pause after 'It is not.'\n\nA camera looks at a street. How many cars just drove past? ... Sounds trivial. It is not.\n\nHello! We are Vinh, Ali and Rexhep. This is our ML4CE semester project, Topic 2.\n\nTwo parts: first we build a car detector ourselves, from scratch. Then we use it to count traffic in a video.\n\nOne thing up front: we show you the numbers that worked - and the ones that did not.\n\n(Straight on to slide 2.)")

# ============================================================ 02 · The task
s = new()
content_header(s, "The task", "Two parts, one pipeline",
               "Detect vehicles in images, then count them in a video")
chip_flow(s, Inches(2.40), [
    ("IMAGES", "1001 dashcam frames"),
    ("DETECTOR", "our own head"),
    ("VIDEO", "60 s of an intersection"),
    ("TRACKER", "IoU + Hungarian, ours"),
    ("COUNT", "per direction"),
])
rows(s, ML, 4.00, [
    ("01", "Part 1 - build a detector by hand",
     "Per cell of a 16 x 16 grid: one objectness score, four box values. Then NMS."),
    ("02", "Part 2 - count vehicles in a video",
     "Fine-tuned YOLO-nano, our tracker for stable IDs, one counting line."),
], gap=0.90)
_, tf = tb(s, ML, Inches(5.90), Inches(11.6), Inches(0.4))
runs(tf.paragraphs[0], [("GRADED ON: WHAT WE TRIED  ·  WHAT WORKED  ·  WHAT DID NOT  ·  "
                         "WHICH METHOD IS BEST, AND WHY", F_MONO, Pt(10.5), SLATE, False, 120)])
footer(s)
notes(s, '~45 SEC.\n\nQuick context: the task sheet asks for two things.\n\nPart 1: a detector. Image in, boxes out. Not a library - our own head on a pretrained backbone. The point is to understand what a detector actually does inside.\n\nPart 2: counting. Same street, but in a video. And that is a different problem, because the count depends on three stages, not one: detect, track, count.\n\nAnd the grade is not only the final number. It is: what did you try, what worked, what did not - and why.')

# ============================================================ 03 · The data
s = new()
content_header(s, "The data", "1001 images - and one video",
               "Kaggle car-object-detection, 676 x 380 px, one vehicle class")
pic_fit(s, "fig_dataset_strip.png", ML, Inches(2.15), CW, Inches(1.80))
caption(s, "Two labelled frames, two frames of empty road - the empty ones are kept as negatives",
        ML, Inches(4.05), w=Inches(9))
FACTS = [("1001", "IMAGES"), ("559", "BOUNDING BOXES"), ("64.5%", "FRAMES WITH NO VEHICLE"),
         ("1.59%", "MEDIAN BOX AREA"), ("2 / 559", "BOXES LOST TO THE GRID")]
xx = 0.667
for value, label in FACTS:
    _, tf = tb(s, Inches(xx), Inches(4.45), Inches(2.4), Inches(0.8))
    runs(tf.paragraphs[0], [(value, F_BLACK, Pt(26), INK, False, 0)], space_after=Pt(2))
    p = tf.add_paragraph()
    runs(p, [(label, F_MONO, Pt(9), SLATE, False, 140)])
    xx += 2.42
_, tf = tb(s, ML, Inches(5.75), Inches(11.6), Inches(0.7))
runs(tf.paragraphs[0], [
    ("All 1001 images are frames of a single video", F_REG, Pt(15), INK, True, 0),
    (", about 0.67 s apart - so neighbours are near-duplicates.", F_LIGHT, Pt(15), INK, False, 0),
], line_spacing=1.4)
footer(s)
notes(s, "~50 SEC.\n\n1001 images, 559 boxes, 676 by 380 pixels. One class: vehicle.\n\nTwo thirds of the frames show empty road. We keep them - they teach the model what 'no car' looks like.\n\nThe cars are small: half of them cover less than 1.6 percent of the image. Remember that number, it comes back later.\n\nAnd now the sentence that decides everything after it: all 1001 images are frames of one single video. Two thirds of a second apart. So two neighbouring images are almost the same picture.\n\n(Say the last line slowly.)")

# ============================================================ 04 · The split
s = new()
content_header(s, "Finding 1 · Methodology", "A random split leaks",
               "Same model, same hyperparameters - only the split changed")
pic_fit(s, "fig_split.png", ML, Inches(2.15), Inches(6.2), Inches(4.1), align="left")
col_label(s, Inches(7.4), Inches(2.30), "Why it happens")
_, tf = tb(s, Inches(7.4), Inches(2.70), Inches(5.2), Inches(2.4))
runs(tf.paragraphs[0], [("A random split puts a frame in train and its twin 0.67 s later in "
                         "test.", F_LIGHT, Pt(14), INK, False, 0)], line_spacing=1.35)
p = tf.add_paragraph()
p.space_before = Pt(14)
runs(p, [("The model is scored on cars it has already memorised.", F_REG, Pt(14), INK, True, 0)],
     line_spacing=1.35)
p = tf.add_paragraph()
p.space_before = Pt(14)
runs(p, [("We split on time instead: 80 / 15 / 5, earliest to latest.",
          F_LIGHT, Pt(14), INK, False, 0)], line_spacing=1.35)
_, tf = tb(s, Inches(7.4), Inches(5.95), Inches(5.2), Inches(0.6))
runs(tf.paragraphs[0], [("EVERY NUMBER FROM HERE ON IS THE TEMPORAL SPLIT",
                         F_MONO, Pt(10.5), SLATE, False, 130)])
footer(s)
notes(s, '~60 SEC.\n\nSo what happens if you split that randomly? One frame goes into training, its twin goes into the test set. Then you are not testing the model any more - you are testing its memory.\n\nSame model, same settings, only the split changed: validation F1 reads 0.78 instead of 0.55. A difference of 0.23 - and it is pure illusion.\n\nWe split along time instead: earliest 80 percent train, next 15 validation, last 5 test. Nothing else changed.\n\nEverything you see from here on is the honest number.\n\n[short pause]')

# ============================================================ 05 · The head
s = new()
content_header(s, "Part 1 · Architecture", "What the head predicts",
               "One objectness score and four box values, per cell of a 16 x 16 grid")
pic_fit(s, "fig_grid.png", ML, Inches(2.10), Inches(4.30), Inches(4.30), align="left")
caption(s, "Network input 512 x 512. Bright cell holds the box centre; the pale cells are our "
           "multi-cell rule", ML, Inches(6.45), w=Inches(4.7))
col_label(s, Inches(5.7), Inches(2.15), "Per cell")
CELL = [("objectness", "is there a vehicle centre in this cell?"),
        ("off_x, off_y", "where in the cell the centre sits, 0 to 1"),
        ("w, h", "box size as a fraction of the image")]
yy = 2.60
for key, body in CELL:
    _, tf = tb(s, Inches(5.7), Inches(yy), Inches(6.9), Inches(0.6))
    runs(tf.paragraphs[0], [(key, F_MONO, Pt(12.5), INK, True, 40)], space_after=Pt(2))
    p = tf.add_paragraph()
    runs(p, [(body, F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.15)
    yy += 0.78
_, tf = tb(s, Inches(5.7), Inches(5.05), Inches(6.9), Inches(1.1))
runs(tf.paragraphs[0], [("All five values live in 0 to 1", F_REG, Pt(14), INK, True, 0),
                        (" - a plain sigmoid produces them, no anchors needed.",
                         F_LIGHT, Pt(14), INK, False, 0)], line_spacing=1.35)
p = tf.add_paragraph()
p.space_before = Pt(8)
runs(p, [("A unit test asserts encode - decode returns the original boxes at IoU 1.0.",
          F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.3)
_, tf = tb(s, Inches(5.7), Inches(6.40), Inches(6.9), Inches(0.4))
runs(tf.paragraphs[0], [("STRIDE 32  ·  16 x 16 GRID  ·  1280 NUMBERS PER IMAGE",
                         F_MONO, Pt(10), SLATE, False, 120)])
footer(s)
notes(s, "~55 SEC.\n\nSo how does the detector work? A network has a fixed output size. But an image can contain zero cars, or seven. You cannot simply 'output a list'.\n\nThe trick: cut the image into a grid, 16 by 16. And ask every cell one small question - is there a car centre inside you? If yes: where exactly, and how big?\n\nThat is five numbers per cell. All between 0 and 1, so a simple sigmoid can produce them. No anchor boxes needed.\n\nAnd one detail that saved us: encoding and decoding are exact opposites, and a test proves it. If those two disagree, the loss still looks perfectly healthy and every box is quietly wrong.")

# ============================================================ 06 · Baseline
s = new()
content_header(s, "Part 1 · Baseline", "The configuration the task sheet prescribes",
               "Frozen ResNet18, one positive cell per box, L1 box loss, 40 epochs")
table(s, ML, 2.30,
      [("Split", 2.2), ("Precision", 1.2), ("Recall", 1.2), ("F1", 1.2), ("AP50", 1.2),
       ("Notes", 3.0)],
      [["Validation (150 images)", "0.679", "0.559", "0.613", "-", "threshold and NMS tuned here"],
       ["Test (50 images)", "0.455", "0.395", "0.423", "0.213", "38 boxes only"]],
      bold_rows=(1,))
pic_fit(s, "fig_curves.png", Inches(6.9), Inches(3.85), Inches(5.7), Inches(2.55))
_, tf = tb(s, ML, Inches(4.05), Inches(6.0), Inches(1.6))
runs(tf.paragraphs[0], [("The test split is too small to carry the argument alone.",
                         F_REG, Pt(14), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(8)
runs(p, [("50 images, 38 boxes - one bad frame moves F1 by several points. The 80/15/5 ratio is "
          "the task sheet's, so validation is our primary number.",
          F_LIGHT, Pt(12.5), MUT, False, 0)], line_spacing=1.35)
_, tf = tb(s, ML, Inches(5.95), Inches(6.0), Inches(0.4))
runs(tf.paragraphs[0], [("TRAINING STOPS HELPING AT EPOCH 17 - MORE DATA, NOT MORE EPOCHS",
                         F_MONO, Pt(10), SLATE, False, 120)])
footer(s)
notes(s, '~50 SEC.\n\nThis is the configuration the task sheet prescribes: ResNet18 frozen, one positive cell per car. We built exactly that first, so every later gain belongs to one specific change.\n\nF1 0.423. Honestly: not good.\n\nTwo caveats we say out loud. The test set is 50 images with 38 boxes - one bad frame moves F1 by several points, so validation is the number we trust.\n\nAnd the curve on the right: after epoch 17 the validation loss stops improving while the training loss keeps falling. More training cannot help here. Only more data, or more varied data.')

# ============================================================ 07 · Ablations
s = new()
content_header(s, "Part 1 · What we tried", "Twelve configurations, one change at a time",
               "Test F1 at IoU 0.5, temporal split, threshold and NMS tuned per run")
pic_fit(s, "fig_ablation.png", ML, Inches(2.10), Inches(8.2), Inches(4.35), align="left")
col_label(s, Inches(9.1), Inches(2.30), "The gap")
_, tf = tb(s, Inches(9.1), Inches(2.70), Inches(3.5), Inches(2.6))
runs(tf.paragraphs[0], [("0.423", F_BLACK, Pt(30), INK, False, 0)], space_after=Pt(0))
p = tf.add_paragraph()
runs(p, [("TASK-SHEET BASELINE", F_MONO, Pt(9), SLATE, False, 140)])
p = tf.add_paragraph()
p.space_before = Pt(12)
runs(p, [("0.904", F_BLACK, Pt(30), INK, False, 0)], space_after=Pt(0))
p = tf.add_paragraph()
runs(p, [("OUR BEST RUN", F_MONO, Pt(9), SLATE, False, 140)])
p = tf.add_paragraph()
p.space_before = Pt(16)
runs(p, [("AP50 goes 0.213 to 0.871 - four times. Two changes account for almost all of it.",
          F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.3)
footer(s)
notes(s, "~55 SEC.\n\nSo we changed one thing at a time. Twelve runs.\n\nDo not read the whole chart - look at the two black bars. At the bottom the task sheet's configuration, 0.423. At the top ours, 0.904. AP50 goes from 0.213 to 0.871, four times better.\n\nEverything in between differs from the baseline in exactly one thing: another backbone, another loss, another assignment rule.\n\nAnd two of those changes explain almost the whole jump. That is the next slide.")

# ============================================================ 08 · The two levers
s = new()
content_header(s, "Part 1 · What worked", "Two changes, almost all of the gain",
               "Each was measured on its own before they were combined")
col_label(s, ML, Inches(2.15), "Lever 1 · unfreeze the backbone     + 0.42 F1")
_, tf = tb(s, ML, Inches(2.55), Inches(5.5), Inches(1.9))
runs(tf.paragraphs[0], [("The frozen features were the bottleneck, not the head.",
                         F_REG, Pt(13.5), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(8)
runs(p, [("ImageNet is object-centric photography; our vehicles are small, blurred and shot from "
          "a dashcam.", F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.3)
col_label(s, Inches(7.0), Inches(2.15), "Lever 2 · multi-cell assignment     + 0.23 F1")
_, tf = tb(s, Inches(7.0), Inches(2.55), Inches(5.6), Inches(1.9))
runs(tf.paragraphs[0], [("One positive cell per box is 453 signals in the whole training set.",
                         F_REG, Pt(13.5), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(8)
runs(p, [("Neighbouring cells fire anyway but were never taught which box to emit, so they emit "
          "fragments.", F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.3)
pic_fit(s, "fig_assign.png", ML, Inches(4.30), Inches(7.6), Inches(2.35), align="left")
_, tf = tb(s, Inches(8.5), Inches(4.60), Inches(4.1), Inches(1.8))
runs(tf.paragraphs[0], [("Train the centre cell plus its two nearest neighbours on the same box.",
                         F_REG, Pt(13), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(8)
runs(p, [("They become agreeing votes that NMS merges. False positives on test: 18 to 10.",
          F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.3)
footer(s)
notes(s, '~75 SEC.\n\nLever one: unfreeze the backbone. Plus 0.42 F1 on its own.\n\nWhy? The backbone was pretrained on ImageNet - object photos, one big sharp dog in the middle. Our cars are small, motion-blurred, filmed from a moving car. Those features simply do not fit, and no head can repair a bad input. The task sheet says you do not need to train the backbone - that is permission, not prohibition. So we report both versions.\n\nLever two: multi-cell assignment. Plus 0.23.\n\nThe rule says only the centre cell is positive. That gives 453 positive signals in the entire training set. But the neighbouring cells fire anyway - they were simply never told which box to predict, so they emit fragments. One car, two errors: a false alarm and a miss at the same time.\n\nWe train the centre cell plus its two neighbours on the same box. Now they agree with each other, and NMS merges them. False positives on the test set drop from 18 to 10.')

# ============================================================ 09 · What did not work
s = new()
content_header(s, "Part 1 · What did not work", "The negative results, kept on purpose",
               "Each was a reasonable idea before it was measured")
table(s, ML, 2.30,
      [("Attempt", 3.3), ("Result", 1.6), ("Why it failed", 6.1)],
      [["CIoU box loss", "0.476 / 0.597", "unstable gradients on tiny targets; L1 is the easier optimisation"],
       ["Focal loss", "0.421", "built for dense anchor detectors; it starved the head"],
       ["Plain BCE", "0.508", "collapsed to 'no vehicle' - 99% cell accuracy, useless detector"],
       ["Flip + colour jitter", "0.485 vs 0.423", "the data lacks scale and viewpoint, not colour"],
       ["Larger input, 640 / 768 px", "0.917 / 0.889", "finds the same 33 vehicles - one false positive fewer"],
       ["Finer grid, stride 16", "0.822", "resolution bought, semantic depth paid - that is what an FPN solves"]],
      row_h=0.52)
_, tf = tb(s, ML, Inches(6.15), Inches(11.6), Inches(0.4))
runs(tf.paragraphs[0], [("WE GOT THIS TABLE WRONG TWICE OURSELVES: A CHECKPOINT READ "
                         "MID-TRAINING, BOX AREA IN INPUT PIXELS",
                         F_MONO, Pt(9.5), SLATE, False, 120)])
footer(s)
notes(s, '~55 SEC.\n\nAnd these are the ideas that lost. We keep them on the slide on purpose.\n\nCIoU loss - in theory better, because it optimises exactly what we measure. It was worse. Focal loss - the worst run of all. Our augmentation, flip and colour jitter - changed nothing, because this data does not lack colour. It lacks scale and viewpoint.\n\nThe interesting row is the fourth one: 640 pixels has the better F1, but it finds exactly the same 33 cars. One false positive fewer, that is all. So we kept the 512 the task sheet specifies - claiming an improvement on one single box would be dishonest.\n\nAnd the last line: we got this table wrong twice ourselves before it was right.')

# ============================================================ 10 · Where the misses are
s = new()
content_header(s, "Part 1 · Error analysis", "Every miss is a small, distant vehicle",
               "Precision holds above 0.95 out to recall 0.87, then the curve falls off a cliff")
pic_fit(s, "fig_pr.png", ML, Inches(2.10), Inches(5.2), Inches(4.3), align="left")
pic_fit(s, "fig_recall_size.png", Inches(6.3), Inches(2.30), Inches(6.3), Inches(3.4))
_, tf = tb(s, Inches(6.3), Inches(5.80), Inches(6.3), Inches(0.8))
runs(tf.paragraphs[0], [("Above 1% of image area, recall is 22 / 22.", F_REG, Pt(14), INK, True, 0)],
     line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(5)
runs(p, [("The fix is a feature pyramid or more data - not more pixels.",
          F_LIGHT, Pt(12.5), MUT, False, 0)], line_spacing=1.3)
footer(s)
notes(s, '~50 SEC.\n\nSo where are the misses? On the left: precision stays above 0.95 up to recall 0.87 - and then the curve falls off a cliff. About 13 percent of the cars are never found, at any threshold.\n\nOn the right: recall by car size. Above one percent of the image we find every single one, 22 out of 22. Below that we start losing them.\n\nSo every miss is a small, distant car. That reads like a resolution problem - which is why we tested bigger inputs. It is not one. The real fix is a feature pyramid, or more data.')

# ============================================================ 11 · Part 1 result
s = new()
content_header(s, "Part 1 · Result", "Test precision 0.943, recall 0.868",
               "MobileNetV3, multi-cell assignment, backbone fine-tuned - 512 px as specified")
STATS = [("0.943", "PRECISION"), ("0.868", "RECALL"), ("0.904", "F1"), ("0.871", "AP50"),
         ("33 / 2 / 5", "TP / FP / FN")]
xx = 0.667
for value, label in STATS:
    _, tf = tb(s, Inches(xx), Inches(2.20), Inches(2.4), Inches(0.9))
    runs(tf.paragraphs[0], [(value, F_BLACK, Pt(34), INK, False, 0)], space_after=Pt(2))
    p = tf.add_paragraph()
    runs(p, [(label, F_MONO, Pt(9.5), SLATE, False, 140)])
    xx += 2.42
pic_fit(s, "fig_part1_examples.png", ML, Inches(3.45), CW, Inches(2.45))
caption(s, "Test frames. Grey = ground truth, black = prediction with its confidence",
        ML, Inches(6.05), w=Inches(8))
_, tf = tb(s, Inches(8.6), Inches(6.00), Inches(4.0), Inches(0.6))
runs(tf.paragraphs[0], [("Precision is a lower bound", F_REG, Pt(11.5), INK, True, 0),
                        (" - some frames leave visible cars unboxed.",
                         F_LIGHT, Pt(11.5), MUT, False, 0)], line_spacing=1.25)
footer(s)
notes(s, '~45 SEC.\n\nSo this is Part 1: precision 0.943, recall 0.868, F1 0.904. Measured on frames from the end of the video that the model has never seen.\n\nGrey is the ground truth, black is our prediction - they sit right on top of each other.\n\nAnd one honest caveat: some frames label one car and leave other clearly visible cars unlabelled. A correct detection there counts against us as a false alarm. So our precision is a lower bound, not the true value.\n\n(Hand over to Part 2 here.)')

# ============================================================ 12 · Part 2 pipeline
s = new()
content_header(s, "Part 2 · The pipeline", "Detect, track, count",
               "Three stages - the count depends on all three, not just the detector")
chip_flow(s, Inches(2.15), [
    ("FRAME", "1798 frames, 60 s"),
    ("YOLO11n", "fine-tuned, one class"),
    ("TRACKER", "IoU + Hungarian, ours"),
    ("LINE COUNTER", "centre crosses, once per ID"),
], chip_h=Inches(0.85))
pic_fit(s, "pipeline_frame660.jpg", ML, Inches(3.25), Inches(7.6), Inches(3.20), align="left")
col_label(s, Inches(8.5), Inches(3.30), "What you see")
SEEN = [("Green box + #ID", "a confirmed track"),
        ("Orange box", "already counted"),
        ("Red line", "the counting line, y = 0.65"),
        ("Panel", "running count per direction")]
yy = 3.70
for key, body in SEEN:
    _, tf = tb(s, Inches(8.5), Inches(yy), Inches(4.1), Inches(0.5))
    runs(tf.paragraphs[0], [(key, F_MONO, Pt(10.5), INK, True, 40)], space_after=Pt(1))
    p = tf.add_paragraph()
    runs(p, [(body, F_LIGHT, Pt(11), MUT, False, 0)], line_spacing=1.15)
    yy += 0.60
_, tf = tb(s, Inches(8.5), Inches(6.15), Inches(4.1), Inches(0.5))
runs(tf.paragraphs[0], [("Counted when the centre crosses", F_REG, Pt(11.5), INK, True, 0),
                        (" - not when it appears.", F_LIGHT, Pt(11.5), MUT, False, 0)],
     line_spacing=1.25)
footer(s)
notes(s, '~40 SEC.\n\nPart 2. Same street, new problem.\n\nThree stages: the detector finds boxes in every frame, the tracker gives them an identity, and the counter watches one line.\n\nGreen box with a number is a tracked car. Orange means already counted - it cannot count a second time. The red line is the counting line.\n\nAnd the rule, this one matters: a car is counted when its centre crosses the line. Not when it appears. The manual count uses exactly the same rule, otherwise human and machine answer different questions.')

# ============================================================ 13 · Detector choice
s = new()
content_header(s, "Part 2 · Comparing detectors", "Twice the detections, fewer vehicles counted",
               "Same frame, same tracker, same line - only the detector changes")
pic_fit(s, "fig_stock_vs_finetune.png", ML, Inches(2.02), CW, Inches(2.78))
caption(s, "A typical frame: the off-the-shelf model adds parked cars far from the line",
        ML, Inches(4.86), w=Inches(9))
pic_fit(s, "fig_detections.png", ML, Inches(5.15), Inches(7.4), Inches(1.45), align="left")
_, tf = tb(s, Inches(8.4), Inches(5.20), Inches(4.2), Inches(1.4))
runs(tf.paragraphs[0], [("Fine-tuning helped - we predicted it would hurt.",
                         F_REG, Pt(13), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(6)
runs(p, [("Unstable boxes fragment tracks: 21.0 tracks per counted vehicle against 8.7.",
          F_LIGHT, Pt(11.5), MUT, False, 0)], line_spacing=1.3)
footer(s)
notes(s, '~55 SEC.\n\nFirst question: which detector? We fine-tuned YOLO11n on our own data and compared it against the off-the-shelf COCO model.\n\nLook at the two frames - same moment in the video. Ours finds 5 cars, the standard model finds 9.\n\nAnd now the surprise: the model with twice as many detections counts fewer cars. 29 against 47.\n\nWhy? It spends its detections on parked cars and pedestrians that never cross the line. And its boxes jump around, so tracks break apart - 21 tracks per counted car instead of 8.7.\n\nMore detections is not better. And by the way, we predicted the opposite: we thought fine-tuning on dashcam images would hurt on a static street camera.')

# ============================================================ 14 · The tracker
s = new()
content_header(s, "Part 2 · The tracker", "IoU association, written from scratch",
               "Hungarian against a greedy baseline, both unit-tested on synthetic boxes")
rows(s, ML, 2.25, [
    ("01", "Associate detections to tracks by IoU, above 0.3",
     "Hungarian is optimal over the whole frame; greedy takes the best pair first."),
    ("02", "Confirm before counting, forget after 0.33 s unseen",
     "Thresholds in seconds, scaled by the video's own frame rate."),
    ("03", "Count when the centre crosses, once per ID",
     "A proper segment intersection, so the line's infinite extension does not count."),
], gap=0.85, w=Inches(6.1))
pic_fit(s, "trajectories_all.jpg", Inches(7.2), Inches(2.25), Inches(5.4), Inches(3.05))
caption(s, "Every tracked path in the clip - 408 tracks for 47 counted vehicles",
        Inches(7.2), Inches(5.36), w=Inches(5.4))
_, tf = tb(s, ML, Inches(5.30), Inches(6.2), Inches(1.3))
runs(tf.paragraphs[0], [("A bug the unit tests caught, and the video never would have.",
                         F_REG, Pt(13), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(6)
runs(p, [("A centre landing exactly on the line was a third state, so the crossing vanished. On "
          "screen that looks like nothing at all - just a count that is too low.",
          F_LIGHT, Pt(11.5), MUT, False, 0)], line_spacing=1.3)
footer(s)
notes(s, "~55 SEC.\n\nThe tracker is ours, not a library.\n\nIt matches boxes between frames by overlap - a car barely moves in one thirtieth of a second, so its box overlaps itself. Hungarian matching solves that optimally for the whole frame, greedy simply takes the best pair first. On this video: exactly the same result, 47 to 47.\n\nThe thresholds are in seconds, not frames. Sounds like a detail - but 10 frames means 0.33 seconds at 30 fps and 0.17 at 60. On someone else's video the same setting would quietly behave differently.\n\nAnd my favourite: a car whose centre landed exactly on the line was treated as a third state, and the crossing disappeared. On the video that looks like nothing at all - just a count that is too low. A unit test found it, the video never would have.")

# ============================================================ 15 · Choosing the video
s = new()
content_header(s, "Part 2 · The hard part", "Choosing the video, not writing the tracker",
               "The binding requirement is invisible in a still frame")
table(s, ML, 2.22,
      [("Candidate", 3.4), ("det/frame", 1.4), ("Moving tracks", 1.6), ("Best line", 1.6),
       ("Verdict", 4.0)],
      [["Dual carriageway, dusk", "18.7", "8", "3 crossings", "too dense to hand-count"],
       ["US freeway", "16.2", "-", "-", "licence forbids redistribution"],
       ["Daylight, free-flowing", "7.4", "22", "12 / 0", "one direction only"],
       ["T-junction", "9.2", "34", "15 (14 one way)", "looked perfect - traffic disperses"],
       ["Intersection (chosen)", "8.8", "93", "42 = 31 / 11", "one line, both flows"]],
      row_h=0.40, bold_rows=(4,))
pic_fit(s, "trajectories_counted.jpg", Inches(7.4), Inches(4.78), Inches(5.2), Inches(1.80))
_, tf = tb(s, ML, Inches(4.90), Inches(6.4), Inches(1.6))
runs(tf.paragraphs[0], [("\"Both directions are visible\" is not the requirement.",
                         F_REG, Pt(13.5), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(6)
runs(p, [("One line has to exist that both flows cross - and that is only visible in the "
          "trajectories.", F_LIGHT, Pt(12), MUT, False, 0)], line_spacing=1.3)
footer(s)
notes(s, "~50 SEC.\n\nHonestly, the hardest part of Part 2 was not the tracker. It was choosing the video.\n\nFour things have to be true at the same time: static camera, both directions, countable by a human, daylight.\n\nWe picked the T-junction clip. Static, daylight, both directions visible - every check passed. Then we tracked it and it fell apart: the traffic disperses, only 15 of 34 moving cars cross any line at all, and 14 of them the same way.\n\nThe requirement was never 'both directions are visible'. It is: one line exists that both flows cross. And you cannot see that in a still image - only in the trajectories.")

# ============================================================ 16 · The count
s = new()
content_header(s, "Part 2 · The count", "47 counted against a manual 43",
               "Same 60 s, same counting rule for the human and the machine")
pic_fit(s, "fig_counts.png", ML, Inches(2.15), Inches(8.0), Inches(4.0), align="left")
col_label(s, Inches(8.9), Inches(2.25), "Result")
_, tf = tb(s, Inches(8.9), Inches(2.62), Inches(3.7), Inches(2.2))
runs(tf.paragraphs[0], [("+9.3%", F_BLACK, Pt(36), INK, False, 0)], space_after=Pt(2))
p = tf.add_paragraph()
runs(p, [("NET OVER-COUNT, 4 VEHICLES", F_MONO, Pt(9.5), SLATE, False, 140)])
p = tf.add_paragraph()
p.space_before = Pt(12)
runs(p, [("Hungarian and greedy tie: 47 and 47, 29 and 29.", F_REG, Pt(12.5), INK, True, 0)],
     line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(8)
runs(p, [("Without a Kalman filter a dropped frame ends a track, the vehicle returns with a new "
          "ID and crosses twice - so the error goes up, as theory says.",
          F_LIGHT, Pt(11.5), MUT, False, 0)], line_spacing=1.3)
_, tf = tb(s, ML, Inches(6.28), Inches(11.6), Inches(0.4))
runs(tf.paragraphs[0], [("STILL OPEN: THE MANUAL COUNT HAS A TOTAL, NOT A PER-DIRECTION SPLIT",
                         F_MONO, Pt(9.5), SLATE, False, 120)])
footer(s)
notes(s, '~45 SEC.\n\nThe result: 47 cars automatically, 43 counted by hand. Plus four - 9.3 percent over.\n\nThe flat part after 44 seconds is the red light, not a failure.\n\nAnd we over-count for a reason we can name: our tracker has no motion model. One missed frame ends the track, the car comes back with a new ID - and crosses the line a second time.\n\nWhat is still open, and we say it openly: the manual count has a total, but not the split per direction. So we do not claim any per-direction accuracy yet.')

# ============================================================ 17 · The result video
s = new()
eyebrow(s, "Part 2 · Result", y=Inches(0.52))
_, tf = tb(s, Inches(7.3), Inches(0.50), Inches(5.35), Inches(0.3))
runs(tf.paragraphs[0], [("CLICK TO PLAY", F_MONO, Pt(11), SLATE, False, 160)],
     align=PP_ALIGN.RIGHT)
mw, mh = Inches(9.6), Inches(5.4)
mx = Emu(int(ML + (CW - mw) / 2))
s.shapes.add_movie(os.path.join(FIGS, "counted_clip_long.mp4"), mx, Inches(0.95), mw, mh,
                   poster_frame_image=os.path.join(FIGS, "counted_poster_long.png"),
                   mime_type="video/mp4")
_, tf = tb(s, ML, Inches(6.50), Inches(11.6), Inches(0.4))
runs(tf.paragraphs[0], [("30 SECONDS OF THE COUNTED OUTPUT  ·  BOXES, TRACK IDS, THE LINE AND "
                         "THE RUNNING COUNT", F_MONO, Pt(10), SLATE, False, 120)])
footer(s)
notes(s, '~40 SEC. - Play the video, talk over it.\n\n[START VIDEO - let it run]\n\nWatch one car crossing the red line, and the counter stepping up.\n\nAnd look at the orange boxes: those are already counted. They can drive around all day, they will not count again.\n\n(After 10-15 seconds, move on.)')

# ============================================================ 18 · Literature
s = new()
content_header(s, "Part 2 · In context", "9.3% is where a street-level camera lands",
               "Published counting errors, split by camera height")
table(s, ML, 2.30,
      [("System", 4.6), ("Camera", 3.2), ("Count error", 2.0), ("Source", 2.2)],
      [["YOLO + visual rhythm", "top view", "0.85-1.4%", "Ribeiro 2025"],
       ["YOLOv5 + DeepSORT", "CCTV, 15 m up", "1.9%", "Tashkent 2023"],
       ["YOLOv3 + SORT, 6 intersections", "14-40 m", "5.5%", "Khazukov 2020"],
       ["Ours: YOLO11n + IoU/Hungarian, no Kalman", "street level", "9.3%", "this project"],
       ["Best method at low vantage", "under 5 m", "9.9-11.0%", "Pakdamansavoji 2025"],
       ["IoU tracker without Kalman, same ablation", "DOT cameras", "over 37%", "Mandal 2021"]],
      row_h=0.42, bold_rows=(3,))
_, tf = tb(s, ML, Inches(5.45), Inches(11.6), Inches(1.0))
runs(tf.paragraphs[0], [("Camera height explains more than the algorithm does.",
                         F_REG, Pt(14), INK, True, 0)], line_spacing=1.3)
p = tf.add_paragraph()
p.space_before = Pt(6)
runs(p, [("The same footage scores 1.81% top-down against 28.58% in perspective view. On our "
          "clip 43.5% of detections overlap another one.", F_LIGHT, Pt(12.5), MUT, False, 0)],
     line_spacing=1.35)
footer(s)
notes(s, '~25 SEC.\n\nIs 9.3 percent good or bad? That depends on where the camera hangs.\n\nFrom above, published systems reach about one percent. From a low camera like ours: around ten.\n\nSame footage, same detector, only the viewing angle changed: 1.8 percent from above against 28 in perspective view. Camera height explains more than the algorithm does.\n\nSo our number is normal for a street-level camera - and for a tracker without a Kalman filter, the literature would have predicted far worse.')

# ============================================================ 19 · Limits & next
s = new()
content_header(s, "Limits and next steps", "What we would fix first",
               "Measured limits, not a wish list")
col_label(s, ML, Inches(2.15), "Known limits")
LIMITS = [("Occlusion", "43.5% of detections overlap another one"),
          ("Oversized boxes", "a box up to 22.7% of the frame - its centre crosses too early"),
          ("Net hides gross", "+4 could be 6 double counts and 2 misses"),
          ("One clip", "no confidence interval; one vehicle is already 2.3%")]
yy = 2.55
for key, body in LIMITS:
    _, tf = tb(s, ML, Inches(yy), Inches(5.4), Inches(0.7))
    runs(tf.paragraphs[0], [(key.upper(), F_MONO, Pt(10.5), INK, True, 60)], space_after=Pt(1))
    p = tf.add_paragraph()
    runs(p, [(body, F_LIGHT, Pt(11), MUT, False, 0)], line_spacing=1.2)
    yy += 0.64
pic_fit(s, "fig_oversized.png", ML, Inches(5.10), Inches(3.3), Inches(1.45), align="left")
caption(s, "One box spanning a vehicle and the background - its centre crosses the line at the "
           "wrong moment", Inches(4.20), Inches(5.55), w=Inches(2.3))
col_label(s, Inches(7.0), Inches(2.15), "Next, in order of expected payoff")
NEXT = [("01", "A Kalman filter, or DeepSORT",
         "the published ablation of our exact design says the error lives here"),
        ("02", "A feature pyramid for Part 1",
         "every remaining miss is a vehicle under 1% of the image"),
        ("03", "Scale and viewpoint augmentation",
         "flip and jitter changed nothing; the data lacks scale"),
        ("04", "Finish the per-direction manual count",
         "the only external ground truth the project has")]
yy = 2.55
for num, head, body in NEXT:
    _, tf = tb(s, Inches(7.0), Inches(yy), Inches(5.6), Inches(0.7))
    runs(tf.paragraphs[0], [(num + "  ", F_BLACK, Pt(13), INK, False, 0),
                            (head, F_REG, Pt(12.5), INK, True, 0)], space_after=Pt(1))
    p = tf.add_paragraph()
    runs(p, [(body, F_LIGHT, Pt(11), MUT, False, 0)], line_spacing=1.2)
    yy += 0.80
footer(s)
notes(s, '~20 SEC.\n\nWhere does our error actually live? Cars hide each other - 43.5 percent of our detections overlap another one. That is geometry, not code.\n\nNext steps, in order of effect: a Kalman filter or DeepSORT against the broken tracks. Then a feature pyramid for the small cars in Part 1. And finish the manual count per direction.')

# ============================================================ 20 · Closing
s = new()
eyebrow(s, "Thank you · Questions", y=Inches(1.30))
_, tf = tb(s, ML, Inches(1.85), Inches(12.2), Inches(1.5))
runs(tf.paragraphs[0], [("IT RUNS ON DATA", F_BLACK, Pt(58), INK, False, 0)], line_spacing=0.98)
p = tf.add_paragraph()
runs(p, [("WE HAVE NEVER SEEN", F_BLACK, Pt(58), INK, False, 0)], line_spacing=0.98)
CLOSE = [("Part 1, any folder of images",
          "python -m src.part1.predict --images <dir> [--csv ground_truth.csv]"),
         ("Part 2, any video",
          "python -m src.part2.suggest_line --video new.mp4   then   run_count --video new.mp4"),
         ("Everything reproducible",
          "one config file, 33 unit tests, weights and the output video ship with the code")]
yy = 4.20
for head, body in CLOSE:
    _, tf = tb(s, ML, Inches(yy), Inches(11.6), Inches(0.7))
    runs(tf.paragraphs[0], [(head, F_REG, Pt(13), INK, True, 0)], space_after=Pt(2))
    p = tf.add_paragraph()
    runs(p, [(body, F_MONO, Pt(10.5), MUT, False, 40)], line_spacing=1.2)
    yy += 0.74
footer(s)
notes(s, '~15 SEC. + questions.\n\nOne last point, and it matters after submission: the model gets tested on data we have never seen. Both parts already run on new data - one command for a folder of images, one for a new video. No code changes needed.\n\nWeights, the output video and 33 unit tests ship with the code.\n\nThank you! Questions are very welcome.\n\nLIKELY QUESTIONS\n- Why build a detector if YOLO exists? The task asks for it - and we compare: our head 0.904 F1, fine-tuned YOLO 0.984 on the same split.\n- Is 47 the correct number? The manual total is 43, so we are 9.3 percent over. That is net, not gross - the per-direction split is not recorded yet.\n- Why only 60 seconds? Because a human has to count it once, reliably. Longer clips are not hand-countable.\n- Would it work on your video? Yes, one command - but the counting line has to be measured again, that is what suggest_line does.\n- Why no Kalman filter? Not required by the task. And we can name exactly what it would fix: the broken tracks that cause our over-count.')

prs.save(OUT)
print("saved", OUT, f"{os.path.getsize(OUT)/1024/1024:.1f} MB, {len(prs.slides._sldIdLst)} slides")
